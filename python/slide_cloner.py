"""
SlideCloner (§6, §6.2 ТЗ).

Вызывается из PHP-класса CasesBot\\Presentation\\SlideCloner как подпроцесс.
OOXML-«хирургия»: копирует конкретные слайды из исходных .pptx в новый файл на
уровне сырых частей пакета (python-pptx opc/package API) — слайд, его layout,
master, theme и все связанные media копируются как есть (blob), с ремаппингом
id/media/layout-master, без чтения через объектную модель и пересборки (это и
теряет SmartArt/кастомные шрифты/анимации у PHPOffice-подобных библиотек).

Заметки докладчика НЕ копируются: в них могут быть внутренние метки/теги эксперта
(#кейс#, категория:тег) — не для показа клиенту.

Режимы:
  python slide_cloner.py clone <output.pptx>
      stdin: {"slides": [{"source_path": "...", "slide_number": N}, ...]}
      Создаёт output.pptx из перечисленных слайдов в указанном порядке.

  python slide_cloner.py add-title-slide <path.pptx>
      stdin: {"title": "...", "date": "..."}
      Добавляет титульный слайд первым в уже существующий path.pptx (§5.1 P0.5).
"""

import json
import posixpath
import re
import sys

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import PartFactory
from pptx.opc.packuri import PackURI
from pptx.oxml.ns import qn

# Заметки слайда не копируются — служебные пометки эксперта не для клиента.
SKIP_RELTYPES = {RT.NOTES_SLIDE}

# r:* атрибуты внутри XML части, которые нужно перенаправить на новые rId при
# копировании (встроенные картинки, гиперссылки). SmartArt/диаграммы (r:dm, r:lo,
# r:qs, встроенные OLE/xlsx) вне охвата v1 — редкость для кейс-слайдов с текстом+фото.
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
R_ATTRS = ("id", "embed", "link")


def partname_template(partname):
    """"/ppt/slides/slide12.xml" -> "/ppt/slides/slide%d.xml" (шаблон для next_partname)."""
    base = partname.baseURI
    filename = partname.filename
    name_no_ext, ext = posixpath.splitext(filename)
    prefix = re.match(r"[a-zA-Z]+", name_no_ext).group(0)
    return f"{base}/{prefix}%d{ext}"


def remap_r_attrs(element, id_map):
    for attr in R_ATTRS:
        for el in element.xpath(f".//*[@r:{attr}]"):
            key = f"{{{R_NS}}}{attr}"
            old = el.get(key)
            if old in id_map:
                el.set(key, id_map[old])


def reserve_partname(template, used_partnames):
    """Следующее свободное имя части по шаблону ("/ppt/slideLayouts/slideLayout%d.xml") —
    аналог dest_package.next_partname(), но без обращения к dest_package.iter_parts().

    dest_package.next_partname() считает имя "занятым" только если часть с ним достижима
    по графу связей (обход из корня пакета), а copy_part_and_deps добавляет новую часть в
    cache ДО того, как она к чему-либо привязана — relate_to родителя происходит только
    после разворачивания рекурсии (см. докстрок ниже). Из-за этого при копировании master'а,
    у которого несколько layout'ов (не только тот, что использует текущий слайд), второй и
    последующие layout'ы копируются, пока первый ещё не подключён к графу, и next_partname()
    выдаёт для них то же самое имя повторно — на выходе два разных part с одинаковым именем
    в архиве (PowerPoint потом не может прочесть один из них и молча удаляет при открытии).
    used_partnames ведётся вручную и не зависит от текущей связности графа.
    """
    n = 1
    while True:
        candidate = template % n
        if candidate not in used_partnames:
            used_partnames.add(candidate)
            return PackURI(candidate)
        n += 1


def copy_part_and_deps(source_part, dest_package, cache, used_partnames):
    """Копирует source_part и всё, от чего он (рекурсивно) зависит, в dest_package.

    Возвращает соответствующую часть в dest_package. Повторные вызовы для уже
    скопированной части (например, общий layout/master у нескольких слайдов из
    одной презентации) отдают закэшированный результат, а не дублируют.
    """
    cache_key = (id(source_part.package), str(source_part.partname))
    if cache_key in cache:
        return cache[cache_key]

    new_partname = reserve_partname(partname_template(source_part.partname), used_partnames)
    new_part = PartFactory(new_partname, source_part.content_type, dest_package, source_part.blob)
    cache[cache_key] = new_part  # до рекурсии — на случай циклических ссылок

    id_map = {}
    for rId, rel in source_part.rels.items():
        if rel.reltype in SKIP_RELTYPES:
            continue
        if rel.is_external:
            id_map[rId] = new_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            new_target = copy_part_and_deps(rel.target_part, dest_package, cache, used_partnames)
            id_map[rId] = new_part.relate_to(new_target, rel.reltype)

    if id_map and hasattr(new_part, "_element"):
        remap_r_attrs(new_part._element, id_map)

    return new_part


def next_master_id(dest_prs):
    """Аналог python-pptx'ного _next_id для sldId, но для sldMasterId: PowerPoint нумерует
    их начиная примерно с 2147483648; python-pptx сам это не делает (id на sldMasterId
    формально необязателен по схеме), но реальный PowerPoint всегда его пишет и без него
    может счесть файл повреждённым и «починить» — молча, без явной ошибки при чтении."""
    sldMasterIdLst = dest_prs._element.get_or_add_sldMasterIdLst()
    used = [int(el.get('id')) for el in sldMasterIdLst if el.get('id') is not None]
    return max([2147483647] + used) + 1


def register_master(dest_prs, dest_package, source_master_part, cache, registered_masters):
    """Гарантирует, что master (уже скопированный в cache через слайд/layout) числится
    в p:sldMasterIdLst презентации — иначе PowerPoint может считать файл повреждённым,
    даже если master фактически достижим через связи slide -> layout -> master."""
    key = (id(source_master_part.package), str(source_master_part.partname))
    if key in registered_masters:
        return
    registered_masters.add(key)

    dest_master_part = cache[key]
    rId = dest_prs.part.relate_to(dest_master_part, RT.SLIDE_MASTER)
    sldMasterIdLst = dest_prs._element.get_or_add_sldMasterIdLst()
    el = etree.SubElement(sldMasterIdLst, qn("p:sldMasterId"))
    el.set("id", str(next_master_id(dest_prs)))
    el.set(qn("r:id"), rId)


def strip_default_template_parts(dest_prs):
    """Presentation() (пустой шаблон python-pptx) несёт свой собственный дефолтный
    slideMaster (с ~11 стандартными layout'ами Office Theme и темой) и printerSettings —
    ни то, ни другое реальным слайдам не нужно: слайды всегда используют master,
    скопированный из исходного файла (register_master), а не этот дефолтный.

    printerSettings1.bin — это ссылка на настройки конкретного принтера с машины,
    на которой собирался сам python-pptx, а не наш контент; PowerPoint не может её
    прочитать и на открытии молча удаляет ("не удалось прочитать часть контента"),
    отсюда запрос на восстановление. Отвязываем эти relationship'ы от presentation.xml
    до добавления реального контента — тогда неиспользуемые части (master, layouts,
    theme, printerSettings) становятся недостижимы и save() их просто не запишет."""
    part = dest_prs.part
    for rId, rel in list(part.rels.items()):
        if rel.reltype in (RT.SLIDE_MASTER, RT.PRINTER_SETTINGS, RT.THEME):
            part.drop_rel(rId)

    sldMasterIdLst = dest_prs._element.get_or_add_sldMasterIdLst()
    sldMasterIdLst.getparent().remove(sldMasterIdLst)


def clone_slides(slides_spec, output_path):
    dest_prs = Presentation()
    strip_default_template_parts(dest_prs)
    dest_package = dest_prs.part.package
    cache = {}
    used_partnames = {str(p.partname) for p in dest_package.iter_parts()}
    registered_masters = set()
    source_cache = {}
    size_set = False

    for item in slides_spec:
        source_path = item["source_path"]
        slide_number = item["slide_number"]

        if source_path not in source_cache:
            source_cache[source_path] = Presentation(source_path)
        source_prs = source_cache[source_path]

        if not size_set:
            dest_prs.slide_width = source_prs.slide_width
            dest_prs.slide_height = source_prs.slide_height
            size_set = True

        source_slide = source_prs.slides[slide_number - 1]
        source_slide_part = source_slide.part
        source_layout_part = source_slide_part.part_related_by(RT.SLIDE_LAYOUT)
        source_master_part = source_layout_part.part_related_by(RT.SLIDE_MASTER)

        new_slide_part = copy_part_and_deps(source_slide_part, dest_package, cache, used_partnames)
        register_master(dest_prs, dest_package, source_master_part, cache, registered_masters)

        rId = dest_prs.part.relate_to(new_slide_part, RT.SLIDE)
        dest_prs._element.get_or_add_sldIdLst().add_sldId(rId)

    dest_prs.save(output_path)
    return len(slides_spec)


def find_title_layout(prs):
    """Layout с настоящим плейсхолдером-заголовком (обычным или центрированным титульным).

    Не полагаемся на slide_masters[0].slide_layouts[0] — для пустого Presention() это и
    правда всегда "Title Slide", но после clone_slides() в файле уже есть master(ы),
    скопированные из исходных презентаций, и от их состава/порядка layouts мы не зависим.
    """
    from pptx.enum.shapes import PP_PLACEHOLDER

    title_types = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            for placeholder in layout.placeholders:
                if placeholder.placeholder_format.type in title_types:
                    return layout

    return prs.slide_masters[0].slide_layouts[0]  # запасной вариант, если титульного не нашлось


def add_title_slide(path, title, date):
    prs = Presentation(path)
    layout = find_title_layout(prs)
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = date

    # Новый слайд add_slide() добавляет последним в sldIdLst — переносим его на первое место.
    sldIdLst = prs._element.get_or_add_sldIdLst()
    new_sldId = sldIdLst[-1]
    sldIdLst.remove(new_sldId)
    sldIdLst.insert(0, new_sldId)

    prs.save(path)


def main():
    if len(sys.argv) < 3:
        print("usage: slide_cloner.py clone|add-title-slide <path> [...]", file=sys.stderr)
        sys.exit(1)

    mode, path = sys.argv[1], sys.argv[2]
    sys.stdin.reconfigure(encoding="utf-8")
    payload = json.load(sys.stdin)

    try:
        if mode == "clone":
            count = clone_slides(payload["slides"], path)
            result = {"ok": True, "slide_count": count}
        elif mode == "add-title-slide":
            add_title_slide(path, payload["title"], payload["date"])
            result = {"ok": True}
        else:
            result = {"ok": False, "error": f"неизвестный режим: {mode}"}
    except Exception as e:  # индексация/сборка не должна падать без объяснения причины
        result = {"ok": False, "error": f"{type(e).__name__}: {e}"}

    sys.stdout.reconfigure(encoding="utf-8")
    json.dump(result, sys.stdout, ensure_ascii=False)


if __name__ == "__main__":
    main()
