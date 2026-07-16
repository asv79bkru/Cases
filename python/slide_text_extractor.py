"""
SlideTextExtractor (§6, §6.2 ТЗ).

Вызывается из PHP-класса CasesBot\\Presentation\\SlideTextExtractor как подпроцесс.
Читает текст слайдов (заголовок, буллеты, полный текст) и картинки со слайда через
python-pptx. Печатает в stdout JSON-массив:
[{"slide_number": 1, "title": "...", "text": "...", "is_case": true, "images": ["..."]}, ...]

Слайд считается кейсом, только если в заметках к слайду есть метка CASE_MARKER
(#кейс# по умолчанию) — так эксперт помечает слайды с кейсами прямо в PowerPoint:
метка видна только в заметках докладчика, не в самом слайде и не в режиме показа.
Картинки извлекаются только для слайдов-кейсов и только если указана папка вывода.

Использование: python slide_text_extractor.py <путь_к_pptx> [метка] [папка_для_картинок]
"""

import json
import os
import re
import sys

CASE_MARKER = "#кейс#"


def slide_title(slide, texts):
    title_shape = slide.shapes.title
    if title_shape is not None:
        title = title_shape.text_frame.text.strip()
        if title:
            return title

    return texts[0].splitlines()[0] if texts else ""


def notes_text(slide):
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text.strip()


def slug(value):
    return re.sub(r"[^0-9A-Za-zА-Яа-яЁё]+", "_", value).strip("_")


def save_images(slide, path, slide_number, images_dir):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    file_slug = slug(os.path.splitext(os.path.basename(path))[0])
    saved = []

    for shape_index, shape in enumerate(slide.shapes, start=1):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue

        image = shape.image
        filename = f"{file_slug}_slide{slide_number}_{shape_index}.{image.ext}"
        with open(os.path.join(images_dir, filename), "wb") as f:
            f.write(image.blob)
        saved.append(filename)

    return saved


def extract(path, case_marker, images_dir):
    from pptx import Presentation

    presentation = Presentation(path)
    slides = []

    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)

        notes = notes_text(slide)
        is_case = case_marker.lower() in notes.lower()

        images = []
        if is_case and images_dir:
            images = save_images(slide, path, index, images_dir)

        slides.append({
            "slide_number": index,
            "title": slide_title(slide, texts),
            "text": "\n".join(texts),
            "is_case": is_case,
            "images": images,
        })

    return slides


def main():
    if len(sys.argv) not in (2, 3, 4):
        print("usage: slide_text_extractor.py <path.pptx> [case_marker] [images_dir]", file=sys.stderr)
        sys.exit(1)

    case_marker = sys.argv[2] if len(sys.argv) >= 3 else CASE_MARKER
    images_dir = sys.argv[3] if len(sys.argv) == 4 else None

    if images_dir:
        os.makedirs(images_dir, exist_ok=True)

    slides = extract(sys.argv[1], case_marker, images_dir)
    sys.stdout.reconfigure(encoding="utf-8")
    json.dump(slides, sys.stdout, ensure_ascii=False)


if __name__ == "__main__":
    main()
